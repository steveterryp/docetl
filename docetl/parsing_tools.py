import io
import os
from typing import List, Optional

from litellm import transcription


def whisper_speech_to_text(filename: str) -> List[str]:
    """
    Transcribe speech from an audio file to text using Whisper model via litellm.
    If the file is larger than 25 MB, it's split into 10-minute chunks with 30-second overlap.

    Args:
        filename (str): Path to the mp3 or mp4 file.

    Returns:
        List[str]: Transcribed text.
    """
    import os

    file_size = os.path.getsize(filename)
    if file_size > 25 * 1024 * 1024:  # 25 MB in bytes
        from pydub import AudioSegment

        audio = AudioSegment.from_file(filename)
        chunk_length = 10 * 60 * 1000  # 10 minutes in milliseconds
        overlap = 30 * 1000  # 30 seconds in milliseconds

        chunks = []
        for i in range(0, len(audio), chunk_length - overlap):
            chunk = audio[i : i + chunk_length]
            chunks.append(chunk)

        transcriptions = []

        for i, chunk in enumerate(chunks):
            buffer = io.BytesIO()
            buffer.name = f"temp_chunk_{i}_{os.path.basename(filename)}"
            chunk.export(buffer, format="mp3")
            buffer.seek(0)  # Reset buffer position to the beginning

            response = transcription(model="whisper-1", file=buffer)
            transcriptions.append(response.text)

        return transcriptions
    else:
        with open(filename, "rb") as audio_file:
            response = transcription(model="whisper-1", file=audio_file)

        return [response.text]


def xlsx_to_string(
    filename: str,
    orientation: str = "col",
    col_order: Optional[List[str]] = None,
    doc_per_sheet: bool = False,
) -> List[str]:
    """
    Convert an Excel file to a string representation or a list of string representations.

    Args:
        filename (str): Path to the xlsx file.
        orientation (str): Either "row" or "col" for cell arrangement.
        col_order (Optional[List[str]]): List of column names to specify the order.
        doc_per_sheet (bool): If True, return a list of strings, one per sheet.

    Returns:
        List[str]: String representation(s) of the Excel file content.
    """
    import openpyxl

    wb = openpyxl.load_workbook(filename)

    def process_sheet(sheet):
        if col_order:
            headers = [
                col for col in col_order if col in sheet.iter_cols(1, sheet.max_column)
            ]
        else:
            headers = [cell.value for cell in sheet[1]]

        result = []
        if orientation == "col":
            for col_idx, header in enumerate(headers, start=1):
                column = sheet.cell(1, col_idx).column_letter
                column_values = [cell.value for cell in sheet[column][1:]]
                result.append(f"{header}: " + "\n".join(map(str, column_values)))
                result.append("")  # Empty line between columns
        else:  # row
            for row in sheet.iter_rows(min_row=2, values_only=True):
                row_dict = {
                    header: value for header, value in zip(headers, row) if header
                }
                result.append(
                    " | ".join(
                        [f"{header}: {value}" for header, value in row_dict.items()]
                    )
                )

        return "\n".join(result)

    if doc_per_sheet:
        return [process_sheet(sheet) for sheet in wb.worksheets]
    else:
        return [process_sheet(wb.active)]


def txt_to_string(filename: str) -> List[str]:
    """
    Read the content of a text file and return it as a list of strings (only one element).

    Args:
        filename (str): Path to the txt or md file.

    Returns:
        List[str]: Content of the file as a list of strings.
    """
    with open(filename, "r", encoding="utf-8") as file:
        return [file.read()]


def docx_to_string(filename: str) -> List[str]:
    """
    Extract text from a Word document.

    Args:
        filename (str): Path to the docx file.

    Returns:
        List[str]: Extracted text from the document.
    """
    from docx import Document

    doc = Document(filename)
    return ["\n".join([paragraph.text for paragraph in doc.paragraphs])]


def pptx_to_string(filename: str, doc_per_slide: bool = False) -> List[str]:
    """
    Extract text from a PowerPoint presentation.

    Args:
        filename (str): Path to the pptx file.
        doc_per_slide (bool): If True, return each slide as a separate
            document. If False, return the entire presentation as one document.

    Returns:
        List[str]: Extracted text from the presentation. If doc_per_slide
            is True, each string in the list represents a single slide.
            Otherwise, the list contains a single string with all slides'
            content.
    """
    from pptx import Presentation

    prs = Presentation(filename)
    result = []

    for slide in prs.slides:
        slide_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_content.append(shape.text)

        if doc_per_slide:
            result.append("\n".join(slide_content))
        else:
            result.extend(slide_content)

    if not doc_per_slide:
        result = ["\n".join(result)]

    return result


def azure_di_read(
    filename: str,
    use_url: bool = False,
    include_line_numbers: bool = False,
    include_handwritten: bool = False,
    include_font_styles: bool = False,
    include_selection_marks: bool = False,
    doc_per_page: bool = False,
) -> List[str]:
    """
    > Note to developers: We used [this documentation](https://learn.microsoft.com/en-us/azure/ai-services/document-intelligence/how-to-guides/use-sdk-rest-api?view=doc-intel-4.0.0&tabs=windows&pivots=programming-language-python) as a reference.

    This function uses Azure Document Intelligence to extract text from documents.
    To use this function, you need to set up an Azure Document Intelligence resource:

    1. [Create an Azure account](https://azure.microsoft.com/) if you don't have one
    2. Set up a Document Intelligence resource in the [Azure portal](https://portal.azure.com/#create/Microsoft.CognitiveServicesFormRecognizer)
    3. Once created, find the resource's endpoint and key in the Azure portal
    4. Set these as environment variables:
       - DOCUMENTINTELLIGENCE_API_KEY: Your Azure Document Intelligence API key
       - DOCUMENTINTELLIGENCE_ENDPOINT: Your Azure Document Intelligence endpoint URL

    The function will use these credentials to authenticate with the Azure service.
    If the environment variables are not set, the function will raise a ValueError.

    The Azure Document Intelligence client is then initialized with these credentials.
    It sends the document (either as a file or URL) to Azure for analysis.
    The service processes the document and returns structured information about its content.

    This function then extracts the text content from the returned data,
    applying any specified formatting options (like including line numbers or font styles).
    The extracted text is returned as a list of strings, with each string
    representing either a page (if doc_per_page is True) or the entire document.

    Args:
        filename (str): Path to the file to be analyzed or URL of the document if use_url is True.
        use_url (bool, optional): If True, treat filename as a URL. Defaults to False.
        include_line_numbers (bool, optional): If True, include line numbers in the output. Defaults to False.
        include_handwritten (bool, optional): If True, include handwritten text in the output. Defaults to False.
        include_font_styles (bool, optional): If True, include font style information in the output. Defaults to False.
        include_selection_marks (bool, optional): If True, include selection marks in the output. Defaults to False.
        doc_per_page (bool, optional): If True, return each page as a separate document. Defaults to False.

    Returns:
        List[str]: Extracted text from the document. If doc_per_page is True, each string in the list represents
                   a single page. Otherwise, the list contains a single string with all pages' content.

    Raises:
        ValueError: If DOCUMENTINTELLIGENCE_API_KEY or DOCUMENTINTELLIGENCE_ENDPOINT environment variables are not set.
    """
    import os

    from azure.ai.documentintelligence import DocumentIntelligenceClient
    from azure.ai.documentintelligence.models import AnalyzeDocumentRequest
    from azure.core.credentials import AzureKeyCredential

    key = os.getenv("DOCUMENTINTELLIGENCE_API_KEY")
    endpoint = os.getenv("DOCUMENTINTELLIGENCE_ENDPOINT")

    if key is None:
        raise ValueError("DOCUMENTINTELLIGENCE_API_KEY environment variable is not set")
    if endpoint is None:
        raise ValueError(
            "DOCUMENTINTELLIGENCE_ENDPOINT environment variable is not set"
        )

    document_analysis_client = DocumentIntelligenceClient(
        endpoint=endpoint, credential=AzureKeyCredential(key)
    )

    if use_url:
        poller = document_analysis_client.begin_analyze_document(
            "prebuilt-read", AnalyzeDocumentRequest(url_source=filename)
        )
    else:
        with open(filename, "rb") as f:
            poller = document_analysis_client.begin_analyze_document("prebuilt-read", f)

    result = poller.result()

    style_content = []
    content = []

    if result.styles:
        for style in result.styles:
            if style.is_handwritten and include_handwritten:
                handwritten_text = ",".join(
                    [
                        result.content[span.offset : span.offset + span.length]
                        for span in style.spans
                    ]
                )
                style_content.append(f"Handwritten content: {handwritten_text}")

            if style.font_style and include_font_styles:
                styled_text = ",".join(
                    [
                        result.content[span.offset : span.offset + span.length]
                        for span in style.spans
                    ]
                )
                style_content.append(f"'{style.font_style}' font style: {styled_text}")

    for page in result.pages:
        page_content = []

        if page.lines:
            for line_idx, line in enumerate(page.lines):
                if include_line_numbers:
                    page_content.append(f" Line #{line_idx}: {line.content}")
                else:
                    page_content.append(f"{line.content}")

        if page.selection_marks and include_selection_marks:
            # TODO: figure this out
            for selection_mark_idx, selection_mark in enumerate(page.selection_marks):
                page_content.append(
                    f"Selection mark #{selection_mark_idx}: State is '{selection_mark.state}' within bounding polygon "
                    f"'{selection_mark.polygon}' and has a confidence of {selection_mark.confidence}"
                )

        content.append("\n".join(page_content))

    if doc_per_page:
        return style_content + content
    else:
        return [
            "\n\n".join(
                [
                    "\n".join(style_content),
                    "\n\n".join(
                        f"Page {i+1}:\n{page_content}"
                        for i, page_content in enumerate(content)
                    ),
                ]
            )
        ]


# Define a dictionary mapping function names to their corresponding functions
PARSING_TOOLS = {
    "whisper_speech_to_text": whisper_speech_to_text,
    "xlsx_to_string": xlsx_to_string,
    "txt_to_string": txt_to_string,
    "docx_to_string": docx_to_string,
    "pptx_to_string": pptx_to_string,
    "azure_di_read": azure_di_read,
}